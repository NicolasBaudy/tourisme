# -*- coding: utf-8 -*-
import streamlit as st
import streamlit.components.v1 as components

import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
import random
from datetime import datetime, timedelta
from scipy.signal import savgol_filter
# from pydrive.auth import GoogleAuth
# from pydrive.drive import GoogleDrive
# from io import StringIO
import os
# import base64
# from pandas.plotting import table as pd_table
# from pdflatex import PDFLaTeX
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


# POUR LANCER L'INTERFACE EN LOCAL:
#   streamlit run interface.py

# POUR LANCER L'INTERFACE SUR LE WEB, après avoir mis le code sur le dépot 
# https://github.com/Lee-RoyMannier/tourisme
# https://share.streamlit.io/lee-roymannier/tourisme/main/interface.py

st.set_option('deprecation.showPyplotGlobalUse', False)

# On enregistrera le contenu affiché pour le rapport en pdf
CONTENU_GLOBAL = {}

# TODO: 
# Se référer à 2019 pour faire les prévisions. 

### I - LECTURE DES DONNEES 
def lecture_donnees(data):
    # Formatage de l'index en date
    data = data.set_index(data.columns[0])
    data.index = data.index.map(lambda x: datetime.strptime(x, "%Y-%m-%d").date())

    # Formatage des nombres à vigule en flottant
    data = data.applymap(lambda x: float(x.replace(",", ".")))

    return data


def donnees_aleatoires(t0=datetime(2017, 6, 1).date(), nb_semaines=4*53):
    data = pd.DataFrame()
 
    data.index = [t0+timedelta(days=i*7) for i in range(nb_semaines)]
    for pays in ['FR', 'SU', 'EN', 'IT', 'ES']:
        data[pays] = [random.gauss(10, 4) for i in range(nb_semaines)]

    data.index.name = "Paris"

    return data


### II - MISE EN FORME
month_str = {
    1: "janvier" , 2: "février"  , 3: "mars", 
    4: "avril"   , 5: "mai"      , 6: "juin", 
    7: "juillet" , 8: "août"     , 9: "septembre",
    10:"octobre" , 11:"novembre" , 12:"décembre"}


def duree_str(date1, date2):
    """Ecrit l'interval entre deux dates d'une manière intuitive et sans 
    redondances. Les dates en entrée sont au format de la librairie datetime.
    
    Par exemple, si on en en entrée:
    >>> date1 = datetime(2020, 10, 3)
    >>> date2 = datetime(2020, 10, 10)
    
    Alors l'interval entre les deux date s'écrira: 'du 3 au 10 octobre 2020' à
    la place par exemple de l'écriture redondante: 'du 3 octobre 2020 au 10 
    octobre 2020'.
    
    Si cela est nécessaire, les années et les mois sont précisés pour chaque
    date. Par exemple, on écrira: 'du 3 octobre 2020 au 10 septembre 2021'."""
    
    d1 = min(date1, date2)
    d2 = max(date1, date2)
    
    def day_str(j):
        if j==1:
            return "1er"
        else:
            return str(j)
    
    a1, m1, j1 = str(d1.year), month_str[d1.month], day_str(d1.day)
    a2, m2, j2 = str(d2.year), month_str[d2.month], day_str(d2.day)
    
    if a1==a2 and m1==m2:
        return  j1+" au "+j2+" "+m2+" "+a2
    elif a1==a2 and m1!=m2:    
        return  j1+" "+m1+" au "+j2+" "+m2+" "+a2 
    else:
        return  j1+" "+m1+" "+a1+" au "+j2+" "+m2+" "+a2


def arrondie_str(x):
    corps, decimales = str(x).split('.')
    return corps+','+decimales[:2]


### III - CALCULS
def variation(x, delta=timedelta(days=7)):
    t2 = max(x.index)
    t1 = t2-delta
    return (x[t2]-x[t1])/x[t1]


def variations(data, date1, date2, delta=4*timedelta(7)):
    # Variations en pourcentage
    dt = data.index[-1] - data.index[-2]
    var = 100*(data-data.shift(round(delta/dt)))/data.shift(round(delta/dt))

    # Variations pendant delta, pour toutes les dates entre date1 et date2 
    date_min = max(min(data.index), date1-delta)
    date_max = min(max(data.index), date2)
    var = var[(var.index>=date_min) & (var.index<=date_max)]
    
    # double index avec la date de début et de fin 
    #dates_1, dates_2 = var.index-delta, var.index
    #dates_1.name, dates_2.name = "début", "fin"
    #var.index = [dates_1, dates_2]

    return var


def tops3(data, date1, date2):

    def tops(data, date1, date2):
        data = data[(data.index>=date1) & (data.index<=date2)]
        tops = data.mean().sort_values(ascending=False)
        return tops
    
    var = variations(data, date1, date2, delta=4*timedelta(7))

    tops_volume    = tops(data, date1, date2)
    tops_variation = tops(var , date1, date2)
    tops_potentiel = (tops_variation*tops_volume).sort_values(ascending=False)

    tops3 = pd.DataFrame({
        "top volume"      : list(tops_volume.head(3).index),
        "top progression" : list(tops_variation.head(3).index),
        "top potentiel"   : list(tops_potentiel.head(3).index)}).T

    # tops3 = tops3.applymap(lambda x: nom_pays(x)+"("+x+")")
    tops3.columns = ["1er", "2ème", "3ème"]
    
    return tops3


def tops_pays(recapitualif_x_semaines, fichier, str_top_semaine):
    """ Fonction retournant un tableau du top 3 des pays ayant le plus gros
    Volume, d'un top 3 des pays ayant le plus haut top de progression ainsi
    qu'un top 3 des pays ayant le plus de potentiel
    Exemple:
        top Volume       top Progression        Top Potentiel
    0  'FR', 'BE', 'NL'  'CH', 'IT', 'NL'  'IT', 'CH', 'NL'
    
    recapitualif_x_semaines étant le récapitulatif du nombre de semaine
    exemple: recap_desc_2s 
    et le top_semaine étant le nom de la colonne 
    en string
    exemple: "TOP 2 SEMAINES"
    """
    top = {"top Volume": [], "top Progression": [], "Top Potentiel": []} 
    
    recapitualif_x_semaines = recapitualif_x_semaines.sort_index()
    recapitualif_x_semaines.fillna(0, inplace=True)
    variation = (variations(fichier, 1).T).sort_index()
    variation.fillna(0, inplace=True)
    concat_tableau = pd.concat([variation, recapitualif_x_semaines], axis=1)
    top_volume = recapitualif_x_semaines.head(3).index.to_list()
    top_progression = variation.sort_values(by=list(variation.columns), 
                                            ascending=False).head(3).index.to_list()
    
    concat_tableau["potentiel"] = concat_tableau[list(concat_tableau.columns)[0]]*concat_tableau[str_top_semaine]
    top_potentiel = list(concat_tableau.sort_values(by=["potentiel"]).head(3).index)
    
    def nettoyage_str(x):
        """ Fonction qui permet de remplacer les "[" ainsi que les "]"
        pour avoir un tableau identique à celui du pdf du client
        """
        x = str(x)
        if "[" and "]" in x:
            x = x.replace("[", "").replace("]", "")
        return x

    top["top Volume"].append(top_volume)
    top["top Progression"].append(top_progression)
    top["Top Potentiel"].append(top_potentiel)
    colonnes = list(top.keys())
    top_3_pays = pd.DataFrame(top, columns=colonnes)
    
    for nom in colonnes:
        top_3_pays[nom] = top_3_pays[nom].apply(nettoyage_str)    
    
    return top_3_pays


def evolutions_sum_annees(fichier, annee):
    """ Tableau les valeurs brutes des 3 dernieres année."""
    evolution_annee = pd.DataFrame()
    
    for i in range(0,4):
        N = fichier["Semaine"].map(lambda x: x.year) == annee - i
        tableau_annee = fichier[N]
        evolution_annee = pd.concat([evolution_annee, tableau_annee])

    evolution_annee["annee"] = evolution_annee["Semaine"].apply(lambda x: str(x)[:4])
    evolution_annee = evolution_annee.reset_index()
    evolution_annee.drop(["index", "Semaine"], axis=1, inplace=True)
    colonne_voulu = list(evolution_annee.columns)[::-1]
    evolution_annee = evolution_annee[colonne_voulu]
    
    return evolution_annee


def evolutions_mois_annee(fichier, mois, annee):
    """ Tableau les valeurs brutes des 3 dernieres année."""
    tableau_date = pd.DataFrame()
    
    for i in range(0,4):
        mois_map = fichier["Semaine"].map(lambda x: x.month) == mois
        tableau_mois = fichier[mois_map]
        annee_map = tableau_mois["Semaine"].map(lambda x: x.year) == annee-i
        tableau_annee_mois = tableau_mois[annee_map]
        tableau_date = pd.concat([tableau_date, tableau_annee_mois])

    tableau_date["annee"] = tableau_date["Semaine"].apply(lambda x: str(x)[:4])
    tableau_date = tableau_date.reset_index()
    tableau_date.drop(["index", "Semaine"], axis=1, inplace=True)
    colonne_voulu = list(tableau_date.columns)[::-1]
    tableau_date = tableau_date[colonne_voulu]
    
    return tableau_date


def valeurs_brutes_3annees(fichier, mois, annee):
    """ Tableau de la sommes des valeurs des pays en fonction du mois et des 
    3 dernières années à partir de l'argument année de la fonction.
    Exemple pour le mois 2 et l'année 2021 :
               Tahiti (PF)  ...  Nouvelle Caledonie (NC)
        annee               ...                         
        2019         215.0  ...                      0.0
        2020         127.0  ...                     25.0
        2021          76.0  ...                      0.0
    """
    tableau_date = pd.DataFrame()
    
    for i in range(0,4):
        mois_map = fichier["Semaine"].map(lambda x: x.month) == mois
        tableau_mois = fichier[mois_map]
        annee_map = tableau_mois["Semaine"].map(lambda x: x.year) == annee-i
        tableau_annee_mois = tableau_mois[annee_map]
        tableau_date = pd.concat([tableau_date, tableau_annee_mois])
    
    tableau_date["annee"] = tableau_date["Semaine"].apply(lambda x: str(x)[:4])
    tableau_brut = tableau_date.groupby("annee").sum()
    
    return tableau_brut


def valeur_trimestrielle(data, annee):
    annee_map = data["Semaine"].map(lambda x: x.year) == annee
    tableau_annee_mois = data[annee_map]
    tableau_annee_mois[str(annee)] = list(tableau_annee_mois.reset_index().index)
    
    return tableau_annee_mois


def moyenne_trimestrielle(data, annee, top6_trimestre):
    """ Moyennes du trimestre T1 de l'année X par rapport a l'année X-1

    Argument:
        data => dataframe
        annee => dernière année 
        trimestre => le numéro du trimestre (T1 est le trimestre correspondant
                                             au 3 premier mois)
        top6_trimestre => tableau regroupant le top 6 trimestriel
    """
    tableau_moyenne = pd.DataFrame()
    def boucle_mois(annee):
        colonnes = list(top6_trimestre.head(6).index)
        tableau_date = pd.DataFrame()
        
        for i in range(0,3):
            mois_map = data["Semaine"].map(lambda x: x.month) == 1+i
            tableau_mois = data[mois_map]
            annee_map = tableau_mois["Semaine"].map(lambda x: x.year) == annee
            tableau_annee_mois = tableau_mois[annee_map]
            tableau_date = pd.concat([tableau_date, tableau_annee_mois])
        return tableau_date.loc[:,colonnes]
    
    tableau_date = boucle_mois(annee)
    moyenne_last_annee = tableau_date.mean()
    tableau_moyenne[annee] = moyenne_last_annee
    tableau_date2 = boucle_mois(annee-1)
    moyenne_annee2 = tableau_date2.mean()
    tableau_moyenne[annee-1] = moyenne_annee2
    tableau_date3 = boucle_mois(annee-2)
    moyenne_annee3 = tableau_date3.mean()
    tableau_moyenne[annee-2] = moyenne_annee3
    
    return tableau_moyenne


def variation_hebdo(data, periode, top_6_hebdo):
    """ Fonction permettant de récupérer 4 semaines (S / S-1 et S-1 / S-2)
    et de calculer les variations sur une période donnée
    Exemple:
        les 2 semaines a partir du 2021-3-21
    retourne un tableau
    """
    colonnes = list(data.columns)
    data = data[data[colonnes[0]] <= periode]
    variation = variations(data,4)
    variation.fillna(0, inplace=True)
    top_6 = list(top_6_hebdo.head(6).index)
    top_variation_hebdo = variation.loc[:,top_6]
    
    return top_variation_hebdo


def variation_trimestrielle(tableau_moyenne):
    def variation_mois_annee(x,y):
        try:
            var = ((x - y) / y) * 100
        except ZeroDivisionError:
            var = 0
        return var
    
    tableau_variation = pd.DataFrame()
    colonnes = list(tableau_moyenne.columns)
    tableau_variation["T1 "+str(colonnes[0])+"/"+str(colonnes[1])] = tableau_moyenne.apply(lambda x: variation_mois_annee(x[colonnes[0]], x[colonnes[1]]), axis=1)
    tableau_variation["T1 "+str(colonnes[0])+"/"+str(colonnes[2])] = tableau_moyenne.apply(lambda x: variation_mois_annee(x[colonnes[0]], x[colonnes[2]]), axis=1)
    
    return tableau_variation

def variation_mensuel(data, annee, mois, top_6_mensuel):
    """ Fonction permettant de récupérer 4 semaines et de calculer les 
    variations sur une période donnée. Nous allons chercher à construire un 
    tableau de moyenne d'un moix X et d'année Y sous la forme:
        
                         Moy Mai 2020    	Moy Mai 2021	
        RÈunion (RE)	               65    	 113    	
        Guadeloupe (GP)	           44    	 69    	
        Martinique (LC)	           28    	 60    	
    """
    moyenne_region = pd.DataFrame()
    top_region = list(top_6_mensuel.head(6).index)
    top_region.append("Semaine")
    data = data.loc[:,top_region]
    for i in range(0,3):
        tmp = pd.DataFrame()
        # Construction du tableau pour l'année N actuel (2021 pour exemple)
        annee_map = data["Semaine"].map(lambda x: x.year) == annee - i
        tableau_mois = data[annee_map]
        # Construction pour le mois M
        mois_map = tableau_mois["Semaine"].map(lambda x: x.month) == mois
        tableau = tableau_mois[mois_map]
        tmp[str(mois)+" "+str(annee-i)] = tableau.head(4).mean(axis=0)
        moyenne_region = pd.concat([moyenne_region, tmp], axis=1)
        
    annees = list(moyenne_region.columns)
    variation = pd.DataFrame()
   
    def variation(x,y):
        if y!=0:
            return ((x - y) / y) * 100
        else:
            return 0
    
    variation[annees[0]+" / "+annees[1]] = moyenne_region.apply(lambda x: variation(x[annees[0]], x[annees[1]]), axis=1)
    variation[annees[0]+" / "+annees[2]] = moyenne_region.apply(lambda x: variation(x[annees[0]], x[annees[2]]), axis=1)
   
    return variation

### IV - GRAPHQUES

def graph_barres(data, nom_x, nom_y, nom_z, formate_date=True):
    # Mise en forme des données (data) pour pouvoir utiliser seaborne, dans un 
    # tableau à trois colonnes (data_graph). La première est le temps, sous 
    # forme de date, la deuxième est les valeurs (volumes, variations, etc...),
    # la troixième les catégories (pays, région, etc..).
    # Les légendes des axes du dessin sont:
    # légende des catégories -> nom_x
    # légende des valeurs    -> nom_y
    # légende du temps       -> nom_z
    data_graph = pd.DataFrame()
    for pays in list(data.columns):
        df = pd.DataFrame({nom_z: data[pays].index, nom_y: data[pays], nom_x: pays})
        data_graph = data_graph.append(df, ignore_index=True)

    # Lorsque les valeurs sont des volumes, les dates représentent des 
    # semaines. Elles sont mises sous un format plus lisible.
    # Lorsque les valeurs sont des variations, les dates représentent le début
    # de la première semaine de variation 
    if formate_date:
        dt = timedelta(days=6) # temps entre le début et la fin de la semaine 
        data_graph[nom_z] = data_graph[nom_z].apply(lambda t: duree_str(t, t+dt))

    # Les volumes sont ensuite représentés à l'aide de barres.
    # Différentes palettes de couleurs ont été testées:
    # YlGnBu RdBu OrRd PRGn Spectral YlOrBr
    fig, ax = plt.subplots(figsize=(10,6), dpi=300)
    sns.barplot(x=nom_x, y=nom_y, hue=nom_z, data=data_graph,
                palette=sns.color_palette("YlGnBu")[-min(len(data),6):])

    # Les différentes semaines sont données en légende
    ax.legend(loc='lower center', bbox_to_anchor=(0.5, 1.01),
              fancybox=True, shadow=False, ncol=3)

    # Les volumes sont écrits en bleu en haut d'une barre, lorsque la valeur
    # est positive et en bas d'une barre lorsque la valeur est négative.
    for p in ax.patches:
        text = " "+format(p.get_height(), '.1f')+" "
        if "%" in nom_y: text+="% "
        x = p.get_x() + p.get_width() / 2.
        y = p.get_height()
        if y >= 0:
            ax.annotate(text, (x,y), ha='center', va='bottom', size=8, 
                        color='blue', xytext=(0,1), textcoords='offset points',
                        rotation=90)
        else:
            ax.annotate(text, (x,y), ha='center', va='top', size=8, 
                        color='red', xytext=(0,1), textcoords='offset points',
                        rotation=90)
    # Des limites un peu plus larges sont fixées en ordonnées afin d'être 
    # certain que les écritures précédentes ne dépassent du cadre
    ymin, ymax = min(data_graph[nom_y]), max(data_graph[nom_y])
    ax.set_ylim([(ymin-0.2*(ymax-ymin) if ymin < 0 else 0),
                 (ymax+0.2*(ymax-ymin) if ymax > 0 else 0)])

    plt.xticks(rotation=45)

    return fig

def graph_3_ans(data, pays, lissage=False):
    """Lissage avec le filtre de Savitzky-Golay . Il utilise les moindres 
    carrés pour régresser une petite fenêtre de vos données sur un polynôme, 
    puis utilise le polynôme pour estimer le point situé au centre de la 
    fenêtre. Enfin, la fenêtre est décalée d’un point de données et le 
    processus se répète. Cela continue jusqu'à ce que chaque point ait été 
    ajusté de manière optimale par rapport à ses voisins. Cela fonctionne très 
    bien même avec des échantillons bruyants provenant de sources non 
    périodiques et non linéaires."""

    a = max(data.index).year
    j1 = data[data.index >= datetime(a, 1, 1).date()].index[0]
    fig, ax = plt.subplots(figsize=(10,6), dpi=300)
    for i in range(3):
        date1, date2 = datetime(a-i, 1, 1).date(), datetime(a-i, 12, 31).date()
        data_ = data[(data.index>date1) & (data.index<=date2)]
        dates = [j1+int((date-date1).days/7.)*timedelta(days=7) for date in data_.index]
        ligne  = ('o--' if i==0 else '.-')
        ligne2 = ('o:'  if i==0 else '.:')
        c = sns.color_palette("YlGnBu")[-i*2-1]
        y = data_[pays].values
        if lissage:
            ylis = savgol_filter(y, 15, 3)
            ax.plot(dates, ylis, ligne, color=c, label=str(a-i)+u" lissé")
            ax.plot(dates, y, ligne2, color=c, label=str(a-i), alpha=0.3)
        else:
            ax.plot(dates, y, ligne, color=c, label=str(a-i))
    
    # Les différentes semaines sont données en légende
    ax.legend(fancybox=True, shadow=False, ncol=1)
    
    # Des limites pour que l'échelle ne change pas entre le lissage et 
    # l'abscence de lissage 
    ax.set_ylim(0, 1.1*data[pays].max())
    
    ax.set_ylabel("Indice Google Trends – Base 100")
    ax.set_title(pays)
    
    plt.xticks([datetime(a, m+1, 1).date() for m in range(12)], 
           ['Janvier', 'Février', 'Mars', 'Avril', 'Mai', 'Juin',
            'Juillet', 'Août', 'Septembre', 'Octobre', 'Novembre', 'Décembre'],
           rotation=45) 

    return fig


### V - GENERATION D'UN RAPPORT
def rapport_pdf():
    pass


### VI - INTERFACES WEB

def entete():
    txt = u"""Bienvenue à l’observatoire digital des destinations françaises
et européennes de Atout France – powered by BC.
L’observatoire digital de Atout France mesure, par quinzaine, par mois 
et par trimestre, les niveaux d’intérêts d’un marché
dans Google Trends (rubrique « travel  ») d’une sélection de mots clés
génériques et des destinations touristiques françaises par espaces
(littoral, outre-mer, urbain, campagne et montagne) et propose,
chaque trimestre, de comparer ces résultats vs les concurrentes en Europe.
(sauf outre-mer).
    """
    cols = st.beta_columns(2) # number of columns in each row! = 2
    cols[0].image("logo_Atout_France.png", use_column_width=True)
    cols[1].image("logo_Baudy_Co.png", use_column_width=True) 
    #cols[1].image("https://nicolasbaudy.files.wordpress.com/2020/02/cropped-logo-new-2.png")
    st.title("Observatoire digital des destinations")
    st.text(txt)


def introduction():
    txt = """
Mesure de l’intérêt des internautes par marché émetteur pour les destinations
françaises (outre-mer, ville,  montagne, campagne et littoral) comparées
entre elles dans un panel donné – rubrique « travel » de Google Trends. 
Les 6 principaux résultats sont proposés (Top 6 – par facilité de lecture,
les valeurs suivantes sont également disponibles) 
- Périodicité d’analyse  : Deux  fois par mois
- Marchés analysés : Allemagne (DE), Belgique (BE), France (FR),
  Pays-Bas (NL) et Royaume-Uni (UK) """

    st.title("Introduction")
    st.header("1- Analyse des recherches pour les destinations françaises")
    st.text(txt)


def visualisation_tops(data):
    date_1, date_2 = max(data.index) - 4*timedelta(7), max(data.index)
    txt = f"""
Synthèse des classements des 3 pays les plus dynamiques sur la période donnée, par défaut 
les 4 dernières semaines disponibles du {duree_str(date_1,date_2)}, respectivement 
pour le 'top volume', la 'top progession' et le 'top potentiel'. 
 
  - L'indicateur de 'volume' est la moyenne des volumes hebdomadaires constatés sur 
les 4 dernières semaines. Il rend compte du niveau d'activité général, tout en 
minimisant les fluctuations pouvant survenir à l'échelle hebdomadaire.
 
  - L'indicateur de 'progression' est la moyenne sur la période des variations
hebdomadaires en pourcentages. Plus il y a eu de variations hebdomadaires à la
hausse pendant 4 semaines, plus l'indicateur de progression est élévé. 

  - L'indicateur de 'potentiel' est le produit de l'indicateur de volume par 
l'indicateur de progression. Il indique les gains potentiels futurs si la tendance 
à la progression observée est conservée.
"""
    st.title("1 - Tops pays")
    st.text(txt)

    date_1, date_2 = max(data.index) - 4*timedelta(7), max(data.index)
    date1 = st.date_input("début:",value=date_1)
    date2 = st.date_input("fin:",  value=date_2)

    top3 = tops3(data, date1, date2)
    
    st.table(top3)
    ax = plt.subplot(111, frame_on=False) # no visible frame
    ax.xaxis.set_visible(False)  # hide the x axis
    ax.yaxis.set_visible(False)  # hide the y axis
    
    return top3.to_latex()


def visualisation_volumes(data):
    txt = """
Google Trends permet de mesurer, de manière relative, l’évolution des recherches 
des internautes, à partir de mots-clés (sujets ou destinations), avec un indice 100
pour la valeur la plus haute au cours de la période analysée. Le champ d’application 
est restreint au domaine du « travel » (ou catégorie « voyage »). Les résultats ne 
sont pas des valeurs absolues mais se lisent en indices.

La visualisation de cet indice au cours des dernières semaines permet de constater 
les fluctuations et les éventuelles tendances de manière empirique. L'attention est 
mise sur les 2 denières semaines puis sur les 4 dernières semaines. """
    st.title("2. Volume des tendances de recherches des deux et quatre dernières semaines ")
    st.text(txt)

    titre_googletrend = "a - Tendances de recherche des 2 dernières semaines"
    st.header(titre_googletrend)
    table = data.tail(2).applymap(lambda x: "{:.1f}".format(x))
    table.index = table.index.map(lambda t: duree_str(t, t+timedelta(days=6)))
    st.write(table)

    nom_x, nom_y, nom_z = "Pays", "Indice Google Trends – Base 100", "Semaine"
    graph_googletrends = graph_barres(data.tail(2), nom_x, nom_y, nom_z)
    st.pyplot(graph_googletrends)

    titre_tendances = "b - Tendances de recherche des 4 dernières semaines"
    st.header(titre_tendances)
    table = data.tail(4).applymap(lambda x: "{:.1f}".format(x))
    table.index = table.index.map(lambda t: duree_str(t, t+timedelta(days=6)))
    st.write(table)
    graph_tendances = graph_barres(data.tail(4), nom_x, nom_y, nom_z)
    st.pyplot(graph_tendances)
    
    resultats = {titre_tendances: graph_tendances,
                 titre_googletrend: graph_googletrends}
    return resultats 


def visualisation_variations(data):
    date = lambda i: max(data.index) + i*timedelta(7)
    semaine =lambda i: duree_str(data.index[-i], data.index[-i]+timedelta(6))
    periode = lambda i, j: "semaine du "+semaine(i)+" à la semaine du "+semaine(j)
    txt = """
D’une semaine à l’autre, les indices des tendances de recherches de Google Trends 
peuvent fluctuer. Les variations sont mesurées Semaine S vs Semaine S-1 et
Semaine S-1 vs Semaine S-2. Les variations S-1 vs S-2 sont comparées à celles
de S-2 vs S-3. """

    st.title("3 - Variations de l'indice")
    st.text(txt)

    st.header("a - Variations S/S-1 comparées à S-1/S-2")
    #st.text(f"{periode(-1,0)}")
    var = variations(data, date(-1), date(0), delta=timedelta(7)).tail(2)
    table = var.applymap(lambda x: "{:.1f}".format(x))
    #table.index = table.index.map(lambda t: duree_str(t, t+timedelta(days=6)))
    table.index = ["S-1/S-2", "S/S-1"]
    st.write(table)

    nom_x, nom_y, nom_z = "Pays", "Variation de l'indice Google Trends – %", "Semaine"
    st.pyplot(graph_barres(var, nom_x, nom_y, nom_z))

    st.header("b - Variations S-1/S-2 comparées à S-2/S-3")
    #st.text(f"{periode(-2,-1)}")
    var = variations(data, date(-2), date(-1), delta=timedelta(7)).tail(2)
    table = var.applymap(lambda x: "{:.1f}".format(x))
    #table.index = table.index.map(lambda t: duree_str(t, t+timedelta(days=6)))
    table.index = ["S-2/S-3", "S-1/S-2"]
    st.write(table)

    nom_x, nom_y, nom_z = "Pays", "Variation de l'indice Google Trends – %", "Semaine"
    st.pyplot(graph_barres(var, nom_x, nom_y, nom_z))


def interface(CONTENU_GLOBAL):    
    def ordre_alpha(categorie):
        """ Pour faciliter la navigation parmi les fichiers, ces derniers sont
        classés par ordre alphabétique. """
        ordonne = sorted(categorie.items(), key=lambda x: x[0])
        categorie = {}
        for donnee in ordonne:
            categorie[donnee[0]] = donnee[1]
        return categorie
    
    def convertion_nom_pays(code_iso):
        """ Nom en Français d'un pays à partir de son code iso en 2 lettres.
        Retourne par exemple "France" pour "FR" """
        try:
            nom_converti = pays.loc[code_iso]["nom_pays"]
            return nom_converti
        except: 
            return code_iso
    
    # Code iso des pays traduits en noms français courts 
    pays = pd.read_csv("iso-pays.csv", header=None)
    pays = pays[[2,4]]
    pays.columns = ["iso", "nom_pays"]
    pays = pays.set_index("iso")
    
    # Lecture des fichiers des tables d'analyse et de leurs noms respectifs
    data_tourisme_pays = {}
    data_tourisme_generique = {}
    emplacement = os.path.join("data_tourisme")
    dossier = os.listdir(emplacement)
    for donnee_tourisme in dossier:
        try:
            donnees_brut = emplacement + "/" + donnee_tourisme
            analyse = pd.read_csv(donnees_brut, sep=";",
                                  encoding = "ISO-8859-1",
                                  engine='python')
            # Le nom du fichier est décomposé pour former le nom qui sera affiché
            decompose = donnee_tourisme.split("_")
            type_analyse = decompose[1]
            type_analyse = type_analyse.split("-")
            # Les analyses générales
            if type_analyse[0] == "Generique":
                # type_analyse = " ".join(type_analyse[:-1])
                nouv_type_analyse = type_analyse[1]
                data_tourisme_generique[nouv_type_analyse] = analyse
            # Les analyses par pays
            else:
                nom_pays = convertion_nom_pays(decompose[0])
                nouv_type_analyse = type_analyse[1]
                if not nom_pays in data_tourisme_pays.keys():
                    data_tourisme_pays[nom_pays] = {}
                data_tourisme_pays[nom_pays][nouv_type_analyse] = analyse
        except:
            pass
        
    # Réorganisation par ordre alphabétique des données
    data_tourisme_pays = ordre_alpha(data_tourisme_pays)
    for pays in data_tourisme_pays:
        data_tourisme_pays[pays] = ordre_alpha(data_tourisme_pays[pays])
    data_tourisme_generique = ordre_alpha(data_tourisme_generique)
    
    entete()
    
    if st.sidebar.checkbox("Présentation", value=True):
        introduction()
     
    # Sélection du type d'analyse à effectuer
    types_analyse = {"Mots clés génériques  par pays": data_tourisme_generique,
                     "Destinations par pays": data_tourisme_pays}
    txt = "Types d'analyses: " 
    noms_types = list(types_analyse.keys())
    mode = st.sidebar.selectbox(txt, noms_types)
    
   
    ### ANALYSE GLOBALE
    if mode == "Mots clés génériques  par pays":
        # Récupération des noms de tables d'analyse et construction de la 
        # liste déroulante
        noms_analyses = list(types_analyse[mode].keys())
        fichier = st.sidebar.selectbox("Quelle analyse effectuer?", noms_analyses)
        data = lecture_donnees(data_tourisme_generique[fichier])
        try:
            ### 1 - LES TOPS
            if st.sidebar.checkbox("1 - Les tops") and fichier != "None":
                top3 = visualisation_tops(data)
                CONTENU_GLOBAL["Top3"] = top3
            
            ### 2 - LES VOLUMES
            if st.sidebar.checkbox("2 - Les volumes") and fichier != "None":
                # On trace les graphiques et le contenu pour le pdf est complété
                graph_volumes = visualisation_volumes(data)
                for graph in graph_volumes:
                    CONTENU_GLOBAL[graph] = graph_volumes[graph]
                    
            ### 3 - LES VARIATIONS
            if st.sidebar.checkbox("3 - Les variations") and fichier != "None":
                visualisation_variations(data)
    
            # COMMENTAIRE à inclure dans le rapport
            #if st.checkbox("Voulez-vous mettre un commentaire?"):
            #    commentaire_1 = st.text_area("Commentaire", "")
    
        except:
            pass

    ### ANALYSE PAR PAYS
    if mode == "Destinations par pays":
        tous_pays = list(data_tourisme_pays.keys())
        pays_choisi = st.sidebar.selectbox("Quel pays?", tous_pays)
        types_analyse = list(data_tourisme_pays[pays_choisi].keys())
        analyse_pays = st.sidebar.selectbox("Quelle analyse effectuer?",
                                           types_analyse)
        data = lecture_donnees(data_tourisme_pays[pays_choisi][analyse_pays])
        try:
            # Date d'analyse
            txt = "Date d'analyse"
            date2 = st.sidebar.date_input(txt,value=max(data.index))

            # Moyennes des volumes sur 2, 4 et 12 semaines, triés par ordre 
            # décroissant
            moyennes = {}
            for i in [2, 4, 12]:
                date1 = date2-i*timedelta(7)
                moyennes[i] = data[(data.index>date1) & (data.index<=date2)].mean()
                moyennes[i] = moyennes[i].sort_values(ascending=False)
                moyennes[i].name = "TOP "+str(i)+" SEMAINES"


            ### 1 - LES TOPS
            if st.sidebar.checkbox("1- Les tops") and analyse_pays != "None":
                st.title("1 - Les tops tendances de recherche - Base : indice 100")
                txt = f"""
Les valeurs moyennes des tendances de recherche de Google Trends sont classées,
sur des périodes, de respectivement:
    - 2 semaines, du {duree_str(date2- 2*timedelta(7), date2)}
    - 4 semaines, du {duree_str(date2- 4*timedelta(7), date2)}
    - 12 semaines, du {duree_str(date2-12*timedelta(7), date2)}"""
                st.text(txt)

                st.header("a - Le top 6")
                cols, k = st.beta_columns(3), 0
                for i, k in zip([2, 4, 12],[0,1,2]):
                    cols[k].table(moyennes[i].apply(arrondie_str).head(6))

                if st.checkbox("afficher les valeurs suivantes..."):
                    st.header("b - Les valeurs suivantes")
                    cols, k = st.beta_columns(3), 0
                    for i, k in zip([2, 4, 12],[0,1,2]):
                        cols[k].table(moyennes[i].apply(arrondie_str).iloc[7:])

                # # COMMENTAIRE à inclure dans le rapport
                # if st.checkbox("Voulez vous mettre un commentaire ?"):
                #     commentaire_2 = st.text_area("Commentaire", "")
           
            ### 2 - LES VOLUMES
            if st.sidebar.checkbox("2 - Les volumes des 3 dernières années du top 6"):
                st.title("2 - Comparaisons annuelles des tops 6")
                classements = ('Top 2 semaines', 'Top 4 semaines','Top 12 semaines')
                lissage    = st.sidebar.checkbox("Lissage") 
                classement = st.sidebar.radio("Classement: ", classements)

                if classement == 'Top 2 semaines':
                    for zone in moyennes[2].head(6).index:
                        st.pyplot(graph_3_ans(data, zone, lissage))
                if classement == 'Top 4 semaines':
                    for zone in moyennes[4].head(6).index:
                        st.pyplot(graph_3_ans(data, zone, lissage))
                if classement == 'Top 12 semaines':
                    for zone in moyennes[12].head(6).index:
                        st.pyplot(graph_3_ans(data, zone, lissage))


            ### 3 - LES VARIATIONS
            if st.sidebar.checkbox("3 - Les variations du top 6 d'une année sur l'autre"):
                txt = """
Les indices de Google Trends, moyennés au choix sur 2, 4 ou 12 dernières semaines
précédant la date d'analyse, sont comparées aux indices sur les mêmes périodes
des années précedentes."""
            
                st.title("3 - Les variations du top 6 d'une année sur l'autre")
                st.text(txt)
                classements = ('2 semaines', '4 semaines', '12 semaines')
                classement = st.sidebar.radio("Moyennes sur: ", classements)

                def moyennes_annuelles(data, periode=i*timedelta(7)):
                    date1 = date2-periode
                    # 1 an avant la date d'analyse:
                    date4 = date2-52*timedelta(7)
                    date3 = date4-periode
                    # 2 ans avant la date d'analyse
                    date6 = date2-104*timedelta(7)
                    date5 = date6-periode
                    
                    moy12 = data[(data.index>date1) & (data.index<=date2)].mean()
                    moy34 = data[(data.index>date3) & (data.index<=date4)].mean()
                    moy56 = data[(data.index>date5) & (data.index<=date6)].mean()
                    df = pd.concat([moy56, moy34, moy12], axis=1)
                    df.columns = [date6.year, date4.year, date2.year]
                    return df.T

                def variations_annuelles(data, periode=i*timedelta(7)):
                    date1 = date2-periode
                    # 1 an avant la date d'analyse:
                    date4 = date2-52*timedelta(7)
                    date3 = date4-periode
                    # 2 ans avant la date d'analyse
                    date6 = date2-104*timedelta(7)
                    date5 = date6-periode
                    
                    moy12 = data[(data.index>date1) & (data.index<=date2)].mean()
                    moy34 = data[(data.index>date3) & (data.index<=date4)].mean()
                    moy56 = data[(data.index>date5) & (data.index<=date6)].mean()
                    df = pd.concat([(moy12-moy56)/moy56*100,
                                    (moy12-moy34)/moy34*100], axis=1)
                    #df.replace([np.inf, -np.inf], 0, inplace=True)
                    df.columns = [str(date2.year) +" vs "+str(date6.year),
                                  str(date2.year) +" vs "+str(date4.year)]
                    return df.T

                if classement == '2 semaines':
                    date1 = date2 - 2*timedelta(7)
                    zones = list(moyennes[2].head(6).index)
                    moy = moyennes_annuelles(data[zones], 2*timedelta(7))
                    var = variations_annuelles(data[zones], 2*timedelta(7))
                if classement == '4 semaines':
                    date1 = date2 - 4*timedelta(7)
                    zones = list(moyennes[4].head(6).index)
                    moy = moyennes_annuelles(data[zones], 4*timedelta(7))
                    var = variations_annuelles(data[zones], 4*timedelta(7))
                if classement == '12 semaines':
                    date1 = date2 - 12*timedelta(7)
                    zones = list(moyennes[12].head(6).index)
                    moy = moyennes_annuelles(data[zones], 12*timedelta(7))
                    var = variations_annuelles(data[zones], 12*timedelta(7))

                titre_var = "a) Valeurs du " + duree_str(date1, date2)
                titre_var += " comparées aux années précédentes."
                st.header(entete)
                
                st.table(moy.T.applymap(lambda x: "{:.1f}".format(x)))
                nom_x, nom_z = u"Régions", "Annees"
                nom_y = "Moyennes de l'indice Google Trends"
                st.pyplot(graph_barres(moy, nom_x, nom_y, nom_z,
                                       formate_date=False))
                
                st.header("b) Variations en %")
                st.table(var.T.applymap(lambda x: "{:.1f}".format(x)))
                nom_y = "Variation des moyennes de l'indice Google Trends - %"
                st.pyplot(graph_barres(var, nom_x, nom_y, nom_z,
                                       formate_date=False))

        except:
            pass
        
    ### EXPORT POWERPOINT
    def ajout_titre(page, type_analyse="", position=0,
                    titre="Indice hebdomadaire des tendances de recherches"):
        """Placement du titre de page. Selon si on on précise un type d'analyse
        ou pas, il sera composé avec une rubique ou non.
        Le titre prend par défaut la première position, mais il est également 
        possible de le placer en dessous d'un autre contenu, avec une position 
        plus élevée."""
        shape = page.shapes[position]
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        p.margin_left = 0
        run = p.add_run()
        run.text = titre
        if type_analyse != "":
            run.text += ' . Rubrique ' + type_analyse
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(18)
        font.bold = True
        font.italic = None
        font.color.rgb = RGBColor(0x11, 0x55, 0xCC)
        
    def table_ppt(page, data, nb_colonnes, nb_lignes, position=0):
        """Création d'une table"""
        taille = nb_colonnes * 1.5
        x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(taille), Inches(3.5)
        shape = page.shapes.add_table(nb_lignes+1, nb_colonnes, x, y, cx, cy)       
        table = shape.table
        index_col = 0
        for nom_colonne in data.columns:
            table.cell(0, index_col).text = nom_colonne
            index_ligne = 1
            for valeur in data[nom_colonne].tolist():
                table.cell(index_ligne, index_col).text = str(valeur)
                index_ligne += 1
            index_col += 1
            
        def iter_cells(table):
            for row in table.rows:
                for cell in row.cells:
                    yield cell
        
        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(10)
        
    # export_ppt = st.sidebar.button("Générer un PowerPoint")
    export_ppt = False # bouton d'export powerpoint caché pour l'instant
    
    if export_ppt: 
        # Pages d'analyse
        # Générique
        # presente = Presentation()
        # page_titre = presente.slide_layouts[1]
        # slide = presente.slides.add_slide(page_titre)
        # slide.shapes.add_picture("logo_Atout_France.png",
        #                           Inches(1), Inches(3),
        #                           width = Inches(5))
        # slide.shapes.add_picture("logo_Baudy_Co.png",
        #                           Inches(4.5), Inches(3),
        #                           width = Inches(5))
        # titre = slide.shapes.title
        # titre.text = u"""Observatoire digital des destinations
        # Analyse Générique"""
        
        
        # # Page des priorités d'action
        # page_priorite = presente.slides.add_slide(page_titre)
        
        # date_1, date_2 = "", ""
        # colonnes = ["", "Top Volume", "Top Progression", "Top Potentiel"]
        # top_quinzaine = pd.DataFrame(columns=colonnes)
        # for type_analyse in data_tourisme_generique:
        #     analyse = lecture_donnees(data_tourisme_generique[type_analyse])
        #     date_2 = max(analyse.index)
        #     date_1 = date_2 - 2*timedelta(7)
        #     top = tops3(analyse, date_1, date_2)
        #     volume = ",".join(top.loc['top volume'])
        #     progression = ",".join(top.loc['top progression'])
        #     potentiel = ",".join(top.loc['top potentiel'])
        #     top_quinzaine.loc[len(top_quinzaine.index)] = [type_analyse, volume,
        #                                                    progression, potentiel]
        # top_priorite = top_quinzaine[["", "Top Progression"]]
        # top_priorite.columns = ["", "Priorité d'action"]
                    
        # titre = "La quinzaine du " + duree_str(date_1, date_2) + " en quelques mots..."
        # ajout_titre(page_priorite, titre=titre, position=0)
        # table_ppt(page_priorite, top_priorite, top_priorite.shape[1],
        #           top_priorite.shape[0], 1)
        
        # # Page des tops de la quinzaine
        # page_top = presente.slides.add_slide(page_titre)
                    
        # titre = "La quinzaine du " + duree_str(date_1, date_2) + " en quelques mots..."
        # ajout_titre(page_top, titre=titre, position=0)
        # table_ppt(page_top, top_quinzaine, top_quinzaine.shape[1],
        #           top_quinzaine.shape[0], 1)
        
        # # Graphiques d'analyse générale
        # for type_analyse in data_tourisme_generique:
        #     page_analyse = presente.slides.add_slide(page_titre)
        #     left = top = Inches(0)
        #     width = Inches(10.0)
        #     height = Inches(0.2)
        #     shape = page_analyse.shapes.add_shape(
        #         MSO_SHAPE.RECTANGLE, left, top, width, height
        #     )
            
        #     ajout_titre(page_analyse, type_analyse)
        
        #     donnees_propres = lecture_donnees(data_tourisme_generique[type_analyse])
        #     graphiques = visualisation_volumes(donnees_propres)
            
        #     decalage = 0
        #     for type_graphique in graphiques:
        #         try:
        #             graph = graphiques[type_graphique]
        #             nom_graph = str(type_analyse) +" "+ str(type_graphique)+".jpg"
        #             image_graph = graph.savefig(nom_graph, dpi=300)
        #             place_image = page_analyse.shapes.add_picture(nom_graph,
        #                                                   Inches(decalage),
        #                                                   Inches(2),
        #                                                   width=Inches(5))
        #             decalage += 5
        #         except:
        #             pass
        #     # break # Arrêt de boucle pour test
        # presente.save('Rapport analyse generique.pptx')
                
        # Par pays
        for pays in data_tourisme_pays:
            presente_pays = Presentation()
            page_titre = presente_pays.slide_layouts[1]
            page_titre_pays = presente_pays.slides.add_slide(page_titre)
            titre_pays = page_titre_pays.shapes.title
            titre_pays.text = "Analyse par Pays"
            page_titre_pays = presente_pays.slides.add_slide(page_titre)
            nom_pays = page_titre_pays.shapes.title
            nom_pays.text = pays
            
            for type_analyse in data_tourisme_pays[pays]:
                donnees_propres = lecture_donnees(data_tourisme_pays[pays][type_analyse])              
                graphiques = {}
                for destination in donnees_propres.columns:
                    nouv_graph = graph_3_ans(donnees_propres, destination)
                    graphiques[destination] = nouv_graph
                
                decalage_x = 0
                decalage_y = 1.5
                page = 1
                for graph_destination in graphiques:
                    if decalage_x == 0 and decalage_y == 1.5:
                        page_analyse = presente_pays.slides.add_slide(page_titre)
                        
                        left = top = Inches(0)
                        width = Inches(10.0)
                        height = Inches(0.2)
                        shape = page_analyse.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE, left, top, width, height
                        )
                        
                        shape = page_analyse.shapes[0]
                        text_frame = shape.text_frame
                        p = text_frame.paragraphs[0]
                        p.margin_left = 0
                        run = p.add_run()
                        run.text = 'Indice hebdomadaire des tendances de recherches'
                        run.text += ' . Rubrique ' + type_analyse
                        run.text += ' . Marché ' + pays
                        
                        font = run.font
                        font.name = 'Calibri'
                        font.size = Pt(18)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0x11, 0x55, 0xCC)

                        page += 1
                    
                    graph = graphiques[graph_destination]
                    nom_graph = str(pays) + " "
                    nom_graph += str(graph_destination)+".jpg"
                    image_graph = graph.savefig(nom_graph, dpi=250)
                    place_image = page_analyse.shapes.add_picture(nom_graph,
                                                          Inches(decalage_x),
                                                          Inches(decalage_y),
                                                          width=Inches(5))
                    if decalage_x > 0 and decalage_y == 1.5:
                        decalage_x = 0
                        decalage_y += 3
                    elif decalage_y > 1.5 and decalage_x > 0:
                        decalage_x = 0
                        decalage_y = 1.5
                    else:
                        decalage_x += 4.7
            presente_pays.save('Rapport analyse '+pays+'.pptx')
            break # Arrêt de boucle pour test
        

### VI - TESTS UNITAIRES
test = False

if test:
    print("lecture des données:")
    try:
        fichier = "../FR-IT-NL-GB-US-BE-CH-DE-ES_Generique-Avion-Hebdo_20210621_1048.csv"
        data = lecture_donnees(fichier)
    except:
        data = donnees_aleatoires(t0=datetime(2017, 6, 1), nb_semaines=4*53)
    print(data)
    
    print("\ntest d'écriture des noms de pays à patir des codes iso:")
    for x in ['FR', 'BE', 'IT', 'CH', 'NL', 'US', 'GB']:
        print("\tcode iso:", x, "=> nom du pays:", x)

    print("\ntest d'écriture d'une durée:")
    date1 = datetime(2021, 5,  9).date()
    date2 = datetime(2021, 5, 30).date()
    print("\tdu", date1, " au ", date2, ": ", duree_str(date1, date2))


### VII - PROGRAMME PRINCIPAL
interface(CONTENU_GLOBAL)